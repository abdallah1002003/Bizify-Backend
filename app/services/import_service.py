import csv
import io
import uuid
from typing import Optional

import docx
import pypdf
from fastapi import HTTPException, UploadFile, status
from pptx import Presentation
from sqlalchemy.orm import Session

from app.models.document import Document
from app.repositories.document_repo import document_repo


class ImportService:
    """Document import and text extraction workflows."""

    async def process_and_save_document(
        self,
        db: Session,
        file: UploadFile,
        current_user_id: uuid.UUID,
    ) -> Document:
        """Extract text from an uploaded file and persist the document."""
        file_bytes = await file.read()
        extracted_text = ""

        try:
            content_type = file.content_type or ""
            filename = (file.filename or "").lower()

            if content_type == "application/pdf" or filename.endswith(".pdf"):
                extracted_text = self._extract_text_from_pdf(file_bytes)
            elif (
                content_type
                in [
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    "application/msword",
                ]
                or filename.endswith(".docx")
            ):
                extracted_text = self._extract_text_from_docx(file_bytes)
            elif (
                content_type
                == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                or filename.endswith(".pptx")
            ):
                extracted_text = self._extract_text_from_pptx(file_bytes)
            elif content_type == "text/csv" or filename.endswith(".csv"):
                extracted_text = self._extract_text_from_csv(file_bytes)
            elif content_type == "text/plain" or filename.endswith(".txt"):
                extracted_text = file_bytes.decode("utf-8")
            else:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Unsupported file format. Supported: PDF, DOCX, PPTX, CSV, TXT.",
                )
        except HTTPException:
            raise
        except Exception as exc:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Error extracting text from document: {exc}",
            ) from exc

        return document_repo.create(
            db,
            obj_in={
                "filename": file.filename,
                "content_type": file.content_type or "text/plain",
                "extracted_text": extracted_text.strip(),
                "user_id": current_user_id,
            },
        )

    def get_document_for_user(
        self,
        db: Session,
        document_id: uuid.UUID,
        user_id: uuid.UUID,
    ) -> Optional[Document]:
        """Fetch a document if it belongs to the user."""
        return document_repo.get_for_user(db, document_id, user_id)

    def delete_document_for_user(
        self,
        db: Session,
        document_id: uuid.UUID,
        user_id: uuid.UUID,
    ) -> Optional[Document]:
        """Delete a user-owned document."""
        document = self.get_document_for_user(db, document_id, user_id)
        if not document:
            return None

        document_repo.delete_instance(db, db_obj=document)
        return document

    def _extract_text_from_pdf(self, file_bytes: bytes) -> str:
        """Extract text from PDF bytes."""
        pdf_reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        text_pages = []
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_pages.append(page_text)
        return "\n".join(text_pages)

    def _extract_text_from_docx(self, file_bytes: bytes) -> str:
        """Extract text from DOCX bytes."""
        document = docx.Document(io.BytesIO(file_bytes))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)

    def _extract_text_from_pptx(self, file_bytes: bytes) -> str:
        """Extract text from PPTX bytes."""
        presentation = Presentation(io.BytesIO(file_bytes))
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _extract_text_from_csv(self, file_bytes: bytes) -> str:
        """Extract text from CSV bytes."""
        content = file_bytes.decode("utf-8")
        reader = csv.reader(io.StringIO(content))
        rows = []
        for row in reader:
            rows.append(" ".join(row))
        return "\n".join(rows)


import_service = ImportService()
